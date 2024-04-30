from openai import OpenAI
import openai
import streamlit as st
import os, time
from pptx import Presentation
from lam import gen_gemma

key = os.getenv("FOSS_API")
key = os.getenv("PKRD_API")

openai.base_url = "https://api.pawan.krd/v1"

client = OpenAI(api_key=key, base_url="https://api.pawan.krd/v1")
def generate_content(subject):
  prompt = "Generate a PowerPoint presentation on the topic of sustainable energy solutions. Include an introduction slide with a captivating title and overview, followed by slides highlighting key concepts such as solar energy, wind power, and hydroelectricity. Each concept should be presented with bullet points summarizing its advantages and applications. Intersperse the presentation with slides containing graphs and charts illustrating global trends in renewable energy adoption and case studies showcasing successful implementation projects. Conclude the presentation with a summary slide outlining key takeaways and recommendations for transitioning towards a sustainable energy future. Additionally, include slides for a Q&A section at the end of the presentation, providing placeholders for audience inquiries and possible answers.in markdown format."
  completion = client.chat.completions.create(
  model="pai-001",
  messages=[
    {"role": "system", "content": prompt},
    {"role": "user", "content": "create a presentation about" + subject + "(only title in first slide)"}
  ]
)

  return (completion.choices[0].message.content)


def create_presentation(content):
  presentation = Presentation()

  # Add slides to the presentation
  slides = content.split("\n\n")
  for slide_content in slides:
      slide = presentation.slides.add_slide(presentation.slide_layouts[1])
      slide.shapes.title.text = slide_content

  # Save the presentation as a PowerPoint file
  presentation.save("/home/inam/slidefy/presentation.pptx")


st.title("Slides Generator")
subject = st.text_input("Enter the subject of the presentation")

def save_md(content):
  with open("presentation.md", "w") as file:
    file.write(content)

if st.button("Generate"):
  start = time.time()
  st.info("Generating presentation slides...")
  #content = generate_content(subject)
  content = gen_gemma(subject)
  print(content)
  save_md(content)
  st.success("Presentation slides generated successfully!")
  end = time.time()
  st.write(f"Time taken: {end - start:.2f} seconds")
  create_presentation(content)

  #give a new link to view another page
  st.markdown("[View the presentation](http://localhost:8501/presentation)")
  st.markdown("[open presentation](./presentation.md)")



#print(generate_content("different types of cars "))